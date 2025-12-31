import os
import sys
import argparse
from pathlib import Path
import tempfile
import fitz  # PyMuPDF
import concurrent.futures
import imghdr
from queue import Queue
from tqdm import tqdm
from vision_api import process_pdf_page
from dotenv import load_dotenv
from pptx import Presentation

# 加载环境变量
load_dotenv()


def process_single_page(page_data):
    """
    处理单个PDF页面
    
    Args:
        page_data: 包含页面处理所需数据的字典
        
    Returns:
        包含页码和处理结果的元组
    """
    page_num = page_data['page_num']
    page = page_data['page']
    temp_dir = page_data['temp_dir']
    api_key = page_data['api_key']
    total_pages = page_data['total_pages']
    
    # print(f"处理第 {page_num + 1} 页，共 {total_pages} 页...")
    
    # 将页面渲染为图像
    image_path = os.path.join(temp_dir, f"page_{page_num + 1}.png")
    pix = page.get_pixmap()
    pix.save(image_path)
    
    # 调用OpenAI视觉模型处理图像
    page_text = process_pdf_page(image_path, api_key)
    
    # 返回页码和处理结果
    return page_num, f"\n\n{page_text}\n\n"


def process_single_slide(slide_data):
    """
    处理单个PPT幻灯片
    
    Args:
        slide_data: 包含幻灯片处理所需数据的字典
        
    Returns:
        包含幻灯片编号和处理结果的元组
    """
    slide_num = slide_data['slide_num']
    slide = slide_data['slide']
    temp_dir = slide_data['temp_dir']
    api_key = slide_data['api_key']
    total_slides = slide_data['total_slides']
    
    # 将幻灯片渲染为图像
    image_path = os.path.join(temp_dir, f"slide_{slide_num + 1}.png")
    
    # 尝试使用win32com导出幻灯片（仅Windows）
    try:
        import win32com.client
        import pythoncom
        
        # 尝试初始化COM（如果尚未初始化）
        try:
            pythoncom.CoInitialize()
        except Exception:
            pass  # COM可能已经初始化，忽略错误
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        
        presentation = powerpoint.Presentations.Open(os.path.abspath(slide_data['ppt_path']))
        presentation.Slides[slide_num + 1].Export(image_path, "PNG")
        
        presentation.Close()
        powerpoint.Quit()
        
        # 反初始化COM
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
        
    except Exception:
        # 备用方法：使用PIL创建包含幻灯片文本的图像
        from PIL import Image, ImageDraw, ImageFont
        
        img = Image.new('RGB', (960, 720), color='white')
        draw = ImageDraw.Draw(img)
        
        # 获取幻灯片中的文本
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        
        # 在图像上绘制文本
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()
        
        y_offset = 50
        for line in slide_text.split('\n')[:20]:
            draw.text((50, y_offset), line, fill='black', font=font)
            y_offset += 30
        
        img.save(image_path)
    
    # 调用OpenAI视觉模型处理图像
    slide_text = process_pdf_page(image_path, api_key)
    
    # 返回幻灯片编号和处理结果
    return slide_num, f"\n\n{slide_text}\n\n"


def process_image_file(image_path: str, output_path: str|None = None, api_key: str|None = None):
    """
    将图片文件转换为Markdown格式
    
    Args:
        image_path: 图片文件路径
        output_path: 输出的Markdown文件路径，如果为None则使用图片文件名
        api_key: OpenAI API密钥
    """
    # 检查图片文件是否存在
    if not os.path.exists(image_path):
        print(f"错误：图片文件 '{image_path}' 不存在")
        return
    
    # 如果未指定输出路径，则使用图片文件名
    if output_path is None:
        output_path = os.path.splitext(image_path)[0] + ".md"
    
    print(f"处理图片文件: {image_path}...")
    
    # 调用OpenAI视觉模型处理图像
    image_text = process_pdf_page(image_path, api_key)
    
    # 写入Markdown文件
    with open(output_path, "w", encoding="utf-8") as md_file:
        md_file.write(f"\n\n{image_text}\n\n")
    
    print(f"转换完成！Markdown文件已保存到: {output_path}")


def convert_pdf_to_markdown(pdf_path: str, output_path: str|None = None, api_key: str|None = None, max_workers: int|None = None):
    """
    将PDF文件转换为Markdown格式
    
    Args:
        pdf_path: PDF文件路径
        output_path: 输出的Markdown文件路径，如果为None则使用PDF文件名
        api_key: OpenAI API密钥
        max_workers: 最大并发线程数，如果为None则从环境变量获取
    """
    # 如果未指定max_workers，则从环境变量获取，默认为5
    if max_workers is None:
        max_workers = int(os.environ.get("MAX_WORKERS", 5))

    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        print(f"错误：PDF文件 '{pdf_path}' 不存在")
        return

    # 如果未指定输出路径，则使用PDF文件名
    if output_path is None:
        output_path = os.path.splitext(pdf_path)[0] + ".md"
    # 打开PDF文件
    try:
        pdf_document = fitz.open(pdf_path)
    except Exception as e:
        print(f"打开PDF文件时出错: {str(e)}")
        return

    # 创建临时目录存储页面图像
    with tempfile.TemporaryDirectory() as temp_dir:
        total_pages = len(pdf_document)
        page_tasks = Queue()
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            page_tasks.put({
                'page_num': page_num,
                'page': page,
                'temp_dir': temp_dir,
                'api_key': api_key,
                'total_pages': total_pages
            })

        results = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_page = {}
            while not page_tasks.empty() and len(future_to_page) < max_workers:
                task = page_tasks.get()
                future = executor.submit(process_single_page, task)
                future_to_page[future] = task

            processed_success = 0
            with tqdm(total=total_pages, desc="页面处理进度", unit="页") as pbar:
                while processed_success < total_pages:
                    while not page_tasks.empty() and len(future_to_page) < max_workers:
                        task = page_tasks.get()
                        future = executor.submit(process_single_page, task)
                        future_to_page[future] = task

                    if not future_to_page and page_tasks.empty():
                        break

                    for future in concurrent.futures.as_completed(list(future_to_page.keys())):
                        task = future_to_page.pop(future)
                        try:
                            page_num, page_content = future.result()
                            results.append((page_num, page_content))
                            processed_success += 1
                            pbar.update(1)
                        except Exception as e:
                            print(f"处理页面时出错: {str(e)}，正在重试该页面 {task['page_num'] + 1}")
                            page_tasks.put(task)

                        while not page_tasks.empty() and len(future_to_page) < max_workers:
                            next_task = page_tasks.get()
                            next_future = executor.submit(process_single_page, next_task)
                            future_to_page[next_future] = next_task
        
        # 按页码顺序组装Markdown内容
        results.sort(key=lambda x: x[0])  # 按页码排序
        markdown_content = ""
        for _, content in results:
            markdown_content += content

        # 关闭PDF文件
        pdf_document.close()

        # 写入Markdown文件
        with open(output_path, "w", encoding="utf-8") as md_file:
            md_file.write(markdown_content)

        print(f"转换完成！Markdown文件已保存到: {output_path}")


def convert_ppt_to_markdown(ppt_path: str, output_path: str|None = None, api_key: str|None = None, max_workers: int|None = None):
    """
    将PPT/PPTX文件转换为Markdown格式
    
    Args:
        ppt_path: PPT/PPTX文件路径
        output_path: 输出的Markdown文件路径，如果为None则使用PPT文件名
        api_key: OpenAI API密钥
        max_workers: 最大并发线程数，如果为None则从环境变量获取
    """
    # 如果未指定max_workers，则从环境变量获取，默认为5
    if max_workers is None:
        max_workers = int(os.environ.get("MAX_WORKERS", 5))

    # 检查PPT文件是否存在
    if not os.path.exists(ppt_path):
        print(f"错误：PPT文件 '{ppt_path}' 不存在")
        return

    # 如果未指定输出路径，则使用PPT文件名
    if output_path is None:
        output_path = os.path.splitext(ppt_path)[0] + ".md"
    
    # 打开PPT文件
    try:
        presentation = Presentation(ppt_path)
    except Exception as e:
        print(f"打开PPT文件时出错: {str(e)}")
        return

    # 创建临时目录存储幻灯片图像
    with tempfile.TemporaryDirectory() as temp_dir:
        total_slides = len(presentation.slides)
        slide_tasks = Queue()
        
        # 准备所有幻灯片任务
        for slide_num, slide in enumerate(presentation.slides):
            slide_tasks.put({
                'slide_num': slide_num,
                'slide': slide,
                'temp_dir': temp_dir,
                'api_key': api_key,
                'total_slides': total_slides,
                'ppt_path': ppt_path
            })

        results = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_slide = {}
            while not slide_tasks.empty() and len(future_to_slide) < max_workers:
                task = slide_tasks.get()
                future = executor.submit(process_single_slide, task)
                future_to_slide[future] = task

            processed_success = 0
            with tqdm(total=total_slides, desc="幻灯片处理进度", unit="页") as pbar:
                while processed_success < total_slides:
                    while not slide_tasks.empty() and len(future_to_slide) < max_workers:
                        task = slide_tasks.get()
                        future = executor.submit(process_single_slide, task)
                        future_to_slide[future] = task

                    if not future_to_slide and slide_tasks.empty():
                        break

                    for future in concurrent.futures.as_completed(list(future_to_slide.keys())):
                        task = future_to_slide.pop(future)
                        try:
                            slide_num, slide_content = future.result()
                            results.append((slide_num, slide_content))
                            processed_success += 1
                            pbar.update(1)
                        except Exception as e:
                            print(f"处理幻灯片时出错: {str(e)}，正在重试该幻灯片 {task['slide_num'] + 1}")
                            slide_tasks.put(task)

                        while not slide_tasks.empty() and len(future_to_slide) < max_workers:
                            next_task = slide_tasks.get()
                            next_future = executor.submit(process_single_slide, next_task)
                            future_to_slide[next_future] = next_task
        
        # 按幻灯片顺序组装Markdown内容
        results.sort(key=lambda x: x[0])
        markdown_content = ""
        for _, content in results:
            markdown_content += content

        # 写入Markdown文件
        with open(output_path, "w", encoding="utf-8") as md_file:
            md_file.write(markdown_content)

        print(f"转换完成！Markdown文件已保存到: {output_path}")


def process_file(file_path: str, output_path: str|None = None, api_key: str|None = None, max_workers: int|None = None):
    """
    处理单个文件（PDF、PPT或图片）
    
    Args:
        file_path: 文件路径
        output_path: 输出的Markdown文件路径
        api_key: OpenAI API密钥
        max_workers: 最大并发线程数
    """
    # 检查文件是否存在
    if not os.path.exists(file_path):
        print(f"错误：文件 '{file_path}' 不存在")
        return
    
    # 检测文件类型
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # 图片文件扩展名列表
    image_extensions = [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"]
    
    # PPT文件扩展名列表
    ppt_extensions = [".ppt", ".pptx"]
    
    # 根据文件类型调用相应的处理函数
    if file_ext == ".pdf":
        convert_pdf_to_markdown(file_path, output_path, api_key, max_workers)
    elif file_ext in ppt_extensions:
        convert_ppt_to_markdown(file_path, output_path, api_key, max_workers)
    elif file_ext in image_extensions:
        process_image_file(file_path, output_path, api_key)
    else:
        # 尝试通过imghdr检测文件类型
        img_type = imghdr.what(file_path)
        if img_type:
            process_image_file(file_path, output_path, api_key)
        else:
            print(f"错误：不支持的文件类型 '{file_ext}'")


def process_files(file_paths, output_dir=None, api_key=None, max_workers=None):
    """
    批量处理多个文件
    
    Args:
        file_paths: 文件路径列表
        output_dir: 输出目录，如果为None则输出到与输入文件相同的目录
        api_key: OpenAI API密钥
        max_workers: 最大并发线程数
    """
    for file_path in file_paths:
        # 如果指定了输出目录，则在该目录下创建输出文件
        output_path = None
        if output_dir:
            output_filename = os.path.splitext(os.path.basename(file_path))[0] + ".md"
            output_path = os.path.join(output_dir, output_filename)
        
        # 处理单个文件
        process_file(file_path, output_path, api_key, max_workers)


def main():
    # 解析命令行参数
    parser = argparse.ArgumentParser(description="将PDF、PPT文件或图片转换为Markdown格式")
    parser.add_argument("file_paths", nargs="+", help="文件路径，支持PDF、PPT/PPTX和常见图片格式（jpg, png等）")
    parser.add_argument("-o", "--output-dir", help="输出目录，默认与输入文件相同目录")
    parser.add_argument("-k", "--api-key", help="OpenAI API密钥")
    parser.add_argument("-w", "--workers", type=int, help="最大并发线程数，仅对PDF和PPT文件有效")

    args = parser.parse_args()

    # 调用处理函数
    process_files(args.file_paths, args.output_dir, args.api_key, args.workers)


if __name__ == "__main__":
    main()
